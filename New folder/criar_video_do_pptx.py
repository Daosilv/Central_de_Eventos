import os
from pptx import Presentation
from moviepy.editor import ImageClip, CompositeVideoClip, concatenate_videoclips, TextClip
from moviepy.video.compositing.transitions import crossfadein
from moviepy.config import change_settings

# Aponta para o executável do ImageMagick
change_settings({"IMAGEMAGICK_BINARY": r"C:\Program Files\ImageMagick-7.1.2-Q16\magick.exe"})


# --- Configurações ---
ARQUIVO_PPTX = "van dia 2408.pptx"
PASTA_IMAGENS_TEMP = "imagens_extraidas"
ARQUIVO_SAIDA = "meu_video_final.mp4"
DURACAO_POR_IMAGEM = 4  # segundos
TEXTO_FINAL = "DEUS RECOMPENSE A CADA UM"
# OTIMIZAÇÃO: Reduz a resolução para HD (720p) para acelerar o processo.
TAMANHO_VIDEO = (1280, 720) 
DURACAO_TRANSICAO = 1 # em segundos

# --- 1. Extração das Imagens do PPTX ---
print(f"Passo 1: Extraindo imagens de '{ARQUIVO_PPTX}'...")
if not os.path.exists(ARQUIVO_PPTX):
    print(f"ERRO: Arquivo '{ARQUIVO_PPTX}' não encontrado.")
    exit()

if not os.path.exists(PASTA_IMAGENS_TEMP):
    os.makedirs(PASTA_IMAGENS_TEMP)

try:
    prs = Presentation(ARQUIVO_PPTX)
    lista_de_imagens_extraidas = []
    contador_imagem = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                nome_arquivo_imagem = os.path.join(PASTA_IMAGENS_TEMP, f"img_{contador_imagem:03d}.png")
                with open(nome_arquivo_imagem, "wb") as f:
                    f.write(image.blob)
                lista_de_imagens_extraidas.append(nome_arquivo_imagem)
                contador_imagem += 1
    print(f"{len(lista_de_imagens_extraidas)} imagens foram extraídas com sucesso.")
except Exception as e:
    print(f"Ocorreu um erro ao extrair as imagens: {e}")
    exit()

# --- 2. Criação do Vídeo com Transições ---
if not lista_de_imagens_extraidas:
    print("Nenhuma imagem encontrada, o vídeo não pode ser criado.")
    exit()

print("\nPasso 2: Criando vídeo com efeitos de transição (modo rápido)...")

try:
    clips_prontos = []
    for caminho_imagem in lista_de_imagens_extraidas:
        clip = ImageClip(caminho_imagem).set_duration(DURACAO_POR_IMAGEM)
        clip_redimensionado = CompositeVideoClip([clip.set_position("center")], size=TAMANHO_VIDEO)
        clips_prontos.append(clip_redimensionado)

    video_com_transicoes = clips_prontos[0] 
    for i in range(1, len(clips_prontos)):
        clipe_com_efeito = crossfadein(clips_prontos[i], DURACAO_TRANSICAO)
        
        video_com_transicoes = concatenate_videoclips([video_com_transicoes, clipe_com_efeito], 
                                                      padding=-DURACAO_TRANSICAO, 
                                                      method="compose")

    clip_texto = TextClip(
        txt=TEXTO_FINAL,
        fontsize=50, # Reduzido um pouco para a nova resolução
        color='white',
        font='Arial-Bold',
        bg_color='black',
        size=TAMANHO_VIDEO
    ).set_duration(DURACAO_POR_IMAGEM)

    video_final = concatenate_videoclips([video_com_transicoes, clip_texto])
    
    # OTIMIZAÇÃO: Adiciona o parâmetro preset="ultrafast" para priorizar velocidade.
    video_final.write_videofile(ARQUIVO_SAIDA, 
                              codec="libx264", 
                              fps=24, 
                              threads=4, 
                              preset="ultrafast")
    print(f"\nVídeo '{ARQUIVO_SAIDA}' criado com sucesso!")

except Exception as e:
    print(f"\nOcorreu um erro ao criar o vídeo: {e}")

finally:
    # --- 3. Limpeza dos Arquivos Temporários ---
    print("Passo 3: Limpando arquivos temporários...")
    if 'video_final' in locals() and hasattr(video_final, 'close'):
        video_final.close()
    
    for imagem in lista_de_imagens_extraidas:
        if os.path.exists(imagem):
            try:
                os.remove(imagem)
            except Exception as e:
                print(f"  Aviso: não foi possível remover o arquivo {imagem}. Erro: {e}")
    if os.path.exists(PASTA_IMAGENS_TEMP):
        try:
            os.rmdir(PASTA_IMAGENS_TEMP)
        except Exception as e:
            print(f"  Aviso: não foi possível remover a pasta {PASTA_IMAGENS_TEMP}. Erro: {e}")

    print("Limpeza concluída.")
