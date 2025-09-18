import os
from pptx import Presentation
# Importações do MoviePy atualizadas
from moviepy.editor import ImageClip, CompositeVideoClip, concatenate_videoclips, TextClip
# NOVO: Importa a configuração do MoviePy
from moviepy.config import change_settings

# NOVO: Aponta para o executável do ImageMagick
# Verifique se este caminho está correto para a sua instalação.
# O padrão é este, mas a versão pode mudar.
change_settings({"IMAGEMAGICK_BINARY": r"C:\Program Files\ImageMagick-7.1.2-Q16\magick.exe"})


# --- Configurações ---
ARQUIVO_PPTX = "van dia 2408.pptx"
PASTA_IMAGENS_TEMP = "imagens_extraidas"
ARQUIVO_SAIDA = "meu_video_final.mp4"
DURACAO_POR_IMAGEM = 4  # segundos
TEXTO_FINAL = "DEUS RECOMPENSE A CADA UM"
TAMANHO_VIDEO = (1920, 1080) # Full HD

# --- 1. Extração das Imagens do PPTX ---
print(f"Passo 1: Extraindo imagens de '{ARQUIVO_PPTX}'...")
if not os.path.exists(ARQUIVO_PPTX):
    print(f"ERRO: Arquivo '{ARQUIVO_PPTX}' não encontrado.")
    exit()

if not os.path.exists(PASTA_IMAGENS_TEMP):
    os.makedirs(PASTA_IMAGENS_TEMP)

try:
    prs = Presentation(ARQUIVO_PPTX)
    lista_de_imagens_extraidas = []
    contador_imagem = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                nome_arquivo_imagem = os.path.join(PASTA_IMAGENS_TEMP, f"img_{contador_imagem:03d}.png")
                with open(nome_arquivo_imagem, "wb") as f:
                    f.write(image.blob)
                lista_de_imagens_extraidas.append(nome_arquivo_imagem)
                contador_imagem += 1
    print(f"{len(lista_de_imagens_extraidas)} imagens foram extraídas com sucesso.")
except Exception as e:
    print(f"Ocorreu um erro ao extrair as imagens: {e}")
    exit()

# --- 2. Criação do Vídeo com Redimensionamento ---
if not lista_de_imagens_extraidas:
    print("Nenhuma imagem encontrada, o vídeo não pode ser criado.")
    exit()

print("\nPasso 2: Criando e redimensionando clipes de imagem...")

try:
    lista_de_clips = []
    for caminho_imagem in lista_de_imagens_extraidas:
        clip = ImageClip(caminho_imagem).set_duration(DURACAO_POR_IMAGEM)
        clip_redimensionado = CompositeVideoClip([clip.set_position("center")], size=TAMANHO_VIDEO)
        lista_de_clips.append(clip_redimensionado)

    clip_imagens = concatenate_videoclips(lista_de_clips)

    clip_texto = TextClip(
        txt=TEXTO_FINAL,
        fontsize=70,
        color='white',
        font='Arial-Bold',
        bg_color='black',
        size=TAMANHO_VIDEO
    ).set_duration(DURACAO_POR_IMAGEM)

    video_final = concatenate_videoclips([clip_imagens, clip_texto])
    video_final.write_videofile(ARQUIVO_SAIDA, codec="libx264", fps=24)
    print(f"\nVídeo '{ARQUIVO_SAIDA}' criado com sucesso!")

except Exception as e:
    print(f"\nOcorreu um erro ao criar o vídeo: {e}")

finally:
    # --- 3. Limpeza dos Arquivos Temporários ---
    print("Passo 3: Limpando arquivos temporários...")
    if 'video_final' in locals():
        video_final.close()
    elif 'clip_imagens' in locals():
        clip_imagens.close()
    
    for imagem in lista_de_imagens_extraidas:
        if os.path.exists(imagem):
            os.remove(imagem)
    if os.path.exists(PASTA_IMAGENS_TEMP):
        os.rmdir(PASTA_IMAGENS_TEMP)
    print("Limpeza concluída.")
